import os
import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd
import markdown

def parse_document(file_path):
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    elif ext in [".txt", ".md"]:
        return parse_text(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def parse_pdf(path):
    doc = fitz.open(path)
    return [page.get_text() for page in doc]

def parse_docx(path):
    doc = docx.Document(path)
    return [para.text for para in doc.paragraphs if para.text.strip()]

def parse_pptx(path):
    prs = pptx.Presentation(path)
    slides = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        slides.append(text.strip())
    return slides

def parse_csv(path):
    df = pd.read_csv(path)
    return [row.to_json() for _, row in df.iterrows()]

def parse_text(path):
    with open(path, "r", encoding="utf-8") as f:
        content = f.read()
    return content.split("\n\n")
